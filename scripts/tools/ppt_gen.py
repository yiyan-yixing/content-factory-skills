#!/usr/bin/env python3
"""PPT 生成器 — 将 Markdown 内容自动生成 PowerPoint 演示文稿。

用法:
  # 从 Markdown 文件生成
  python3 ppt_gen.py article.md -o deck.pptx

  # 从目录下多个 Markdown 文件生成多页
  python3 ppt_gen.py chapters/ -o deck.pptx --title "标题"

  # 自定义模板
  python3 ppt_gen.py article.md -o deck.pptx --template dark --brand "#1a1a2e,#e9c46a"

依赖: pip install python-pptx
"""
import argparse
import glob
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── 品牌色板（默认） ──
BRAND = {
    "bg_dark": RGBColor(0x1A, 0x1A, 0x2E),
    "gold": RGBColor(0xE9, 0xC4, 0x6A),
    "red": RGBColor(0xE7, 0x6F, 0x51),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_bg": RGBColor(0xF5, 0xF3, 0xEF),
    "dark_text": RGBColor(0x2D, 0x2D, 0x2D),
    "accent_teal": RGBColor(0x45, 0xB2, 0x9D),
}

# ── 模板定义 ──
TEMPLATES = {
    "dark": {
        "slide_bg": BRAND["bg_dark"],
        "title_color": BRAND["gold"],
        "body_color": BRAND["white"],
        "accent": BRAND["gold"],
        "subtitle_color": BRAND["white"],
    },
    "light": {
        "slide_bg": BRAND["light_bg"],
        "title_color": BRAND["bg_dark"],
        "body_color": BRAND["dark_text"],
        "accent": BRAND["red"],
        "subtitle_color": BRAND["dark_text"],
    },
    "tech": {
        "slide_bg": RGBColor(0x0D, 0x11, 0x17),
        "title_color": RGBColor(0x58, 0xA6, 0xFF),
        "body_color": RGBColor(0xC9, 0xD1, 0xD9),
        "accent": RGBColor(0x3F, 0xB9, 0x50),
        "subtitle_color": RGBColor(0x8B, 0x94, 0x9E),
    },
}


def parse_markdown_slides(content: str):
    """将 Markdown 内容解析为 slide 列表。
    每个 H2 (##) 或 H3 (###) 是一个新 slide。
    """
    slides = []
    current = None
    lines = content.splitlines()
    buffer = []

    for line in lines:
        if line.startswith("## ") or line.startswith("### "):
            if current is not None:
                current["body"] = "\n".join(buffer).strip()
                slides.append(current)
            level = 1 if line.startswith("## ") else 2
            current = {"title": re.sub(r"^#{2,3} ", "", line), "level": level, "body": ""}
            buffer = []
        elif current is not None:
            buffer.append(line)

    if current is not None:
        current["body"] = "\n".join(buffer).strip()
        slides.append(current)

    # 如果没有 H2/H3，整篇作为一个 slide
    if not slides:
        slides.append({"title": "", "level": 0, "body": content.strip()})

    return slides


def md_to_paragraphs(text: str):
    """将 Markdown 文本解析为段落列表（支持粗体、列表）"""
    paragraphs = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("- ") or line.startswith("* "):
            paragraphs.append(("bullet", line[2:]))
        elif re.match(r"^\d+\. ", line):
            paragraphs.append(("numbered", re.sub(r"^\d+\. ", "", line)))
        else:
            paragraphs.append(("text", line))
    return paragraphs


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, template, left=0.5, top=0.3, width=9.0, height=0.8):
    """添加标题栏（带左侧装饰线）"""
    # 左侧装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = template["accent"]
    line.line.fill.background()

    # 标题文字
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top), Inches(width - 0.2), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = template["title_color"]
    p.font.bold = True


def add_body_text(slide, text, template, left=0.7, top=1.3, width=8.5, height=5.5):
    """添加正文"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs = md_to_paragraphs(text)
    for i, (ptype, ptext) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = ptext
        p.font.size = Pt(16)
        p.font.color.rgb = template["body_color"]

        if ptype == "bullet":
            p.level = 0
            p.font.size = Pt(15)
        elif ptype == "numbered":
            p.font.size = Pt(15)

        p.space_after = Pt(6)


def add_footer(slide, template, text="— 一言一行 —"):
    """添加底部脚注"""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = template["accent"]
    p.alignment = PP_ALIGN.CENTER


def add_cover_slide(prs, title, subtitle, template):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, template["slide_bg"])

    # 中央标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = template["title_color"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.8), Inches(3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = template["accent"]
    line.line.fill.background()

    # 副标题
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = template["subtitle_color"]
        p2.alignment = PP_ALIGN.CENTER

    # 底部脚注
    add_footer(slide, template)


def add_section_slide(prs, title, template):
    """创建章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, template["slide_bg"])

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = template["title_color"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.3), Inches(4), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = template["accent"]
    line.line.fill.background()

    add_footer(slide, template)


def add_content_slide(prs, slide_data, template):
    """创建内容页"""
    level = slide_data.get("level", 1)

    if level == 1 and slide_data["title"]:
        # H2 → 章节分隔（add_section_slide 内部创建 slide）
        add_section_slide(prs, slide_data["title"], template)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if template["slide_bg"] != BRAND["light_bg"]:
        set_slide_bg(slide, template["slide_bg"])
    else:
        # 浅色模板用白色背景
        set_slide_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))

    add_title_bar(slide, slide_data["title"], template)
    add_body_text(slide, slide_data["body"], template)
    add_footer(slide, template)


def generate_pptx(markdown_paths: list, output: str,
                  title: str = "", subtitle: str = "",
                  template_name: str = "dark",
                  author: str = "一言一行"):
    """从 Markdown 文件列表生成 PPT"""
    template = TEMPLATES.get(template_name, TEMPLATES["dark"])
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 元数据
    prs.core_properties.author = author

    # 封面
    add_cover_slide(prs, title or "演示文稿", subtitle, template)

    # 内容页
    for md_path in markdown_paths:
        content = Path(md_path).read_text(encoding="utf-8")
        slides = parse_markdown_slides(content)
        for slide_data in slides:
            add_content_slide(prs, slide_data, template)

    # 尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, template["slide_bg"])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢"
    p.font.size = Pt(44)
    p.font.color.rgb = template["title_color"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save(output)
    size = Path(output).stat().st_size
    print(f"✅ PPT 已生成: {output} ({size / 1024:.0f} KB)")
    return output


def main():
    parser = argparse.ArgumentParser(description="Markdown → PPT 生成器")
    parser.add_argument("input", nargs="+", help="Markdown 文件或目录")
    parser.add_argument("-o", "--output", default="output.pptx", help="输出 PPTX 路径")
    parser.add_argument("--title", default="", help="演示文稿标题")
    parser.add_argument("--subtitle", default="", help="副标题")
    parser.add_argument("--template", default="dark",
                        choices=list(TEMPLATES.keys()),
                        help="模板 (dark/light/tech)")
    parser.add_argument("--author", default="一言一行", help="作者")

    args = parser.parse_args()

    # 收集所有 markdown 文件
    md_files = []
    for inp in args.input:
        p = Path(inp)
        if p.is_dir():
            md_files.extend(sorted(glob.glob(str(p / "**/*.md"), recursive=True)))
        elif p.suffix in (".md", ".txt"):
            md_files.append(str(p))

    if not md_files:
        print("❌ 错误: 未找到 Markdown 文件")
        sys.exit(1)

    print(f"📄 输入文件: {len(md_files)} 个")
    generate_pptx(md_files, args.output,
                  title=args.title, subtitle=args.subtitle,
                  template_name=args.template, author=args.author)


if __name__ == "__main__":
    main()
